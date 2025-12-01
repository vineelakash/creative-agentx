from pptx import Presentation
from pptx.util import Inches, Pt
import os

def create_presentation(output_path: str):
    """
    Generates a PowerPoint presentation for the Creative Agent project.
    Requires: pip install python-pptx
    """
    prs = Presentation()

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Creative Agent X"
    subtitle.text = "Kaggle Agents Intensive Capstone - Creative Track"

    # Slide 2: Pipeline Architecture
    slide_layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Pipeline Architecture"
    content = slide.placeholders[1]
    content.text = (
        "1. Input: Topic, Audience, Tone\n"
        "2. Generator (LLM/Offline): \n"
        "   - Outline Generation\n"
        "   - Blog Expansion\n"
        "   - Video Script Generation\n"
        "   - Image Prompt Generation\n"
        "3. Evaluator: Quality & Safety Checks\n"
        "4. Output: JSON & Text Files"
    )

    # Slide 3: Example Outputs
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Example Outputs"
    content = slide.placeholders[1]
    content.text = (
        "Topic: Future of Space Travel\n\n"
        "Blog Excerpt:\n"
        "'Since the dawn of time, humanity has looked up at the stars...'\n\n"
        "Video Script:\n"
        "[SCENE START] INT. SPACESHIP..."
    )

    # Slide 4: Evaluation Metrics
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Evaluation Metrics"
    content = slide.placeholders[1]
    content.text = (
        "- Flesch Reading Ease Score\n"
        "- Word Count Verification\n"
        "- Keyword Coverage Analysis\n"
        "- Safety & Toxicity Checks"
    )

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    output_file = os.path.join(os.path.dirname(__file__), "creative_agent_presentation.pptx")
    try:
        create_presentation(output_file)
    except ImportError:
        print("Error: 'python-pptx' library not found. Please run: pip install python-pptx")
    except Exception as e:
        print(f"An error occurred: {e}")
